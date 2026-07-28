#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""知识库文档提取器 —— 把 PDF / DOCX / PPTX / XLSX 提取为可检索的 Markdown。

为什么需要:工作台全文搜索只能读纯文本,二进制文档里的内容检索不到、AI 也引用不到。
本工具把它们提取成 `.extracted/` 下的同构 md 镜像,既可被搜索,也可作为知识编译(wiki)的输入。

用法:
  python3 bin/extract.py                 # 增量提取知识库(仅处理新增/改动)
  python3 bin/extract.py --all           # 全量重提取
  python3 bin/extract.py --dir "路径"     # 指定目录
  python3 bin/extract.py --status        # 只看统计,不提取

依赖(均为本机现成):pdftotext(poppler)/ textutil(macOS)/ python-pptx / openpyxl
"""
import argparse
import hashlib
import json
import re
import subprocess
import sys
from datetime import datetime
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent


def mask_dirs():
    """哪些知识库子目录在提取时需要脱敏 —— **由用户在界面上指定,不由程序推断**。

    默认为空 = 都不脱敏。理由是用户的:知识库里很多客户信息本来就出自公开发表的
    文章,脱敏既没必要、又会把可读的参考资料搅乱。哪些数据该保护,用户自己最清楚。
    """
    try:
        cfg = json.loads((ROOT / "workbench.json").read_text(encoding="utf-8"))
        return [d for d in (cfg.get("knowledge", {}).get("mask_dirs") or []) if d]
    except Exception:      # noqa: BLE001
        return []


def _mask(text, rel, dirs):
    """只对指定目录下的文件脱敏。替换完全对照映射表,表里没有的一律不动。"""
    if not dirs:
        return text
    top = Path(rel).parts[0] if Path(rel).parts else ""
    if top not in dirs:
        return text
    try:
        import mask
        return mask.mask_text(text)
    except Exception:      # noqa: BLE001  映射表不可用时不能让提取整体失败
        return text


DEFAULT_SRC = ROOT / "knowledge" / "my local knowledge"
OUT_DIR = ROOT / "knowledge" / ".extracted"
MANIFEST = OUT_DIR / "_manifest.json"

SUPPORTED = {".pdf", ".docx", ".doc", ".pptx", ".xlsx", ".rtf", ".txt", ".md", ".csv"}
SKIP_DIRS = {".extracted", "reports", "__MACOSX"}
MAX_CHARS = 300_000        # 单文件提取上限,防超大文件拖垮检索


def sha(path, limit=4 * 1024 * 1024):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        h.update(f.read(limit))
    return h.hexdigest()[:16]


def clean(text):
    """压掉提取产生的噪声:连续空行、页眉页脚式重复短行、控制字符。"""
    text = text.replace("\x00", "").replace("\f", "\n\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [ln.rstrip() for ln in text.splitlines()]
    # 出现 >5 次且长度 <25 的短行,多为页眉/页脚/水印
    from collections import Counter
    cnt = Counter(ln.strip() for ln in lines if 0 < len(ln.strip()) < 25)
    noise = {k for k, v in cnt.items() if v > 5}
    out = [ln for ln in lines if ln.strip() not in noise]
    return re.sub(r"\n{3,}", "\n\n", "\n".join(out)).strip()


# ---------------- 各格式提取 ----------------

def from_pdf(p):
    r = subprocess.run(["pdftotext", "-layout", "-enc", "UTF-8", str(p), "-"],
                       capture_output=True, timeout=180)
    return r.stdout.decode("utf-8", "replace") if r.returncode == 0 else ""


def _textutil(p):
    """textutil 底层调用(不做兜底,避免与 from_docx 互相递归)。"""
    r = subprocess.run(["textutil", "-convert", "txt", "-inputencoding", "UTF-8",
                        "-stdout", str(p)], capture_output=True, timeout=120)
    if r.returncode != 0:
        r = subprocess.run(["textutil", "-convert", "txt", "-stdout", str(p)],
                           capture_output=True, timeout=120)
    return r.stdout.decode("utf-8", "replace") if r.returncode == 0 else ""


def from_docx(p):
    """docx 优先 python-docx:能拿到表格,且对 textutil 返回空的文件仍有效(实测)。"""
    try:
        import docx
    except ImportError:
        return _textutil(p)
    d = docx.Document(str(p))
    parts = []
    for para in d.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        lv = 0
        try:
            m = re.match(r"Heading (\d)", para.style.name or "")
            lv = int(m.group(1)) if m else 0
        except (AttributeError, ValueError):
            lv = 0
        parts.append(("#" * min(lv, 4) + " " + t) if lv else t)
    for i, tb in enumerate(d.tables, 1):
        rows = []
        for r in tb.rows:
            cells = [c.text.strip().replace("\n", " ") for c in r.cells]
            if any(cells):
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"\n[表格 {i}]\n" + "\n".join(rows[:200]))
    out = "\n\n".join(parts)
    return out if out.strip() else _textutil(p)


def from_word(p):
    """textutil 支持 doc/rtf/rtfd/odt/html/webarchive/wordml。
    -inputencoding 对无 charset 声明的 html/txt 至关重要,否则中文乱码。"""
    return _textutil(p)


def from_json(p):
    """JSON → 可读文本(保留结构,便于当模板骨架)。"""
    try:
        obj = json.loads(p.read_text(encoding="utf-8", errors="replace"))
        return "```json\n" + json.dumps(obj, ensure_ascii=False, indent=2) + "\n```"
    except (json.JSONDecodeError, OSError):
        return p.read_text(encoding="utf-8", errors="replace")


def from_csv(p):
    """CSV → Markdown 表格。"""
    import csv as _csv
    import io as _io
    txt = p.read_text(encoding="utf-8", errors="replace")
    rows = list(_csv.reader(_io.StringIO(txt)))[:500]
    if not rows:
        return ""
    out = ["| " + " | ".join(rows[0]) + " |",
           "|" + "---|" * len(rows[0])]
    out += ["| " + " | ".join(r + [""] * (len(rows[0]) - len(r))) + " |"
            for r in rows[1:]]
    return "\n".join(out)


def from_pptx(p):
    from pptx import Presentation
    prs = Presentation(str(p))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        buf = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text.strip()
                if t:
                    buf.append(t)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        buf.append(" | ".join(cells))
        if slide.has_notes_slide:
            n = (slide.notes_slide.notes_text_frame.text or "").strip()
            if n:
                buf.append(f"[备注] {n}")
        if buf:
            parts.append(f"## 第 {i} 页\n\n" + "\n\n".join(buf))
    return "\n\n".join(parts)


def from_xlsx(p):
    import openpyxl
    wb = openpyxl.load_workbook(str(p), data_only=True, read_only=True)
    parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c).strip() for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
            if len(rows) > 500:
                rows.append("…(超 500 行,已截断)")
                break
        if rows:
            parts.append(f"## 工作表:{ws.title}\n\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def from_text(p):
    return p.read_text(encoding="utf-8", errors="replace")


EXTRACTORS = {
    ".pdf": from_pdf,
    # textutil 系(macOS 自带)
    ".docx": from_docx, ".doc": from_word, ".rtf": from_word, ".rtfd": from_word,
    ".odt": from_word, ".html": from_word, ".htm": from_word, ".webarchive": from_word,
    ".wordml": from_word,
    # 结构化
    ".pptx": from_pptx, ".xlsx": from_xlsx, ".xlsm": from_xlsx,
    ".csv": from_csv, ".json": from_json,
    # 纯文本
    ".txt": from_text, ".md": from_text, ".markdown": from_text, ".log": from_text,
}


def extract_one(src, src_root):
    suf = src.suffix.lower()
    fn = EXTRACTORS.get(suf)
    if not fn:
        return None, "unsupported"
    try:
        raw = fn(src)
    except Exception as e:  # noqa: BLE001
        return None, f"error: {type(e).__name__}: {e}"[:120]
    body = clean(raw)[:MAX_CHARS]
    if len(body) < 20:
        return None, "empty(可能是扫描件/图片型文档,需 OCR)"
    rel = src.relative_to(src_root)
    head = (f"---\nsource: {rel}\ntype: {suf.lstrip('.')}\n"
            f"extracted: {datetime.now().strftime('%Y-%m-%d %H:%M')}\n"
            f"chars: {len(body)}\n---\n\n# {src.stem}\n\n"
            f"> 由 {src.name} 自动提取(原文见 `knowledge/my local knowledge/{rel}`)\n\n")
    return head + body + "\n", "ok"


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--dir", default=str(DEFAULT_SRC))
    ap.add_argument("--all", action="store_true", help="全量重提取")
    ap.add_argument("--status", action="store_true", help="只统计")
    args = ap.parse_args()

    src_root = Path(args.dir).resolve()
    if not src_root.is_dir():
        print(f"[extract] 源目录不存在: {src_root}")
        sys.exit(1)

    files = [p for p in sorted(src_root.rglob("*"))
             if p.is_file() and not p.name.startswith(".")
             and p.suffix.lower() in SUPPORTED
             and not any(d in p.parts for d in SKIP_DIRS)]

    manifest = {}
    if MANIFEST.exists() and not args.all:
        try:
            manifest = json.loads(MANIFEST.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            manifest = {}

    if args.status:
        done = sum(1 for k in manifest if manifest[k].get("status") == "ok")
        print(f"[extract] 源文档 {len(files)} 份 · 已提取 {done} 份 · "
              f"输出目录 {OUT_DIR.relative_to(ROOT)}")
        for k, v in list(manifest.items()):
            if v.get("status") != "ok":
                print(f"  ⚠ {k}: {v.get('status')}")
        return

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    MDIRS = mask_dirs()
    if MDIRS:
        print(f"[extract] 脱敏目录: {'、'.join(MDIRS)}(其余目录原样提取)")
    stats = {"ok": 0, "skip": 0, "fail": 0}
    for p in files:
        rel = str(p.relative_to(src_root))
        digest = sha(p)
        prev = manifest.get(rel)
        out_path = OUT_DIR / (Path(rel).with_suffix("").as_posix() + ".md")
        if prev and prev.get("sha") == digest and prev.get("status") == "ok" \
                and out_path.exists():
            stats["skip"] += 1
            continue
        body, status = extract_one(p, src_root)
        if body is None:
            manifest[rel] = {"sha": digest, "status": status}
            stats["fail"] += 1
            print(f"  ⚠ {rel} → {status}")
            continue
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_text(_mask(body, rel, MDIRS), encoding="utf-8")
        manifest[rel] = {"sha": digest, "status": "ok",
                         "out": str(out_path.relative_to(OUT_DIR)),
                         "chars": len(body)}
        stats["ok"] += 1
        print(f"  ✓ {rel} → {len(body)} 字")

    MANIFEST.write_text(json.dumps(manifest, ensure_ascii=False, indent=1),
                        encoding="utf-8")
    print(f"[extract] 完成:新提取 {stats['ok']} · 跳过(未变) {stats['skip']} · "
          f"失败 {stats['fail']} · 输出 {OUT_DIR.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
